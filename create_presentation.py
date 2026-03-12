#!/usr/bin/env python3
"""
AI in Software Development - Management Buy-in Presentation
McKinsey-style professional presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import zipfile
import os
from lxml import etree

# Constants
NAVY = RGBColor(0x05, 0x1C, 0x2C)
CYAN = RGBColor(0x00, 0xA9, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MED_GRAY = RGBColor(0x66, 0x66, 0x66)
LINE_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
BG_GRAY = RGBColor(0xF2, 0xF2, 0xF2)

TITLE_SIZE = Pt(22)
BODY_SIZE = Pt(14)
SUB_HEADER_SIZE = Pt(18)
HEADER_SIZE = Pt(28)

LM = Inches(0.8)
CONTENT_W = Inches(11.7)
SW = Inches(13.333)
SH = Inches(7.5)

def _clean_shape(shape):
    sp = shape._element
    style = sp.find(qn('p:style'))
    if style is not None:
        sp.remove(style)

def set_ea_font(run, typeface='KaiTi'):
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', typeface)

def add_text(slide, left, top, width, height, text, font_size=BODY_SIZE,
             font_name='Arial', font_color=DARK_GRAY, bold=False,
             alignment=PP_ALIGN.LEFT, ea_font='KaiTi', anchor=MSO_ANCHOR.TOP,
             line_spacing=Pt(6), bullet=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    anchor_map = {MSO_ANCHOR.MIDDLE: 'ctr', MSO_ANCHOR.BOTTOM: 'b', MSO_ANCHOR.TOP: 't'}
    bodyPr.set('anchor', anchor_map.get(anchor, 't'))
    for attr in ['lIns','tIns','rIns','bIns']:
        bodyPr.set(attr, '45720')
    
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if bullet:
            line = '• ' + line
        p.text = line
        p.font.size = font_size
        p.font.name = font_name
        p.font.color.rgb = font_color
        p.font.bold = bold
        p.alignment = alignment
        p.space_before = line_spacing if i > 0 else Pt(0)
        p.space_after = Pt(0)
        p.line_spacing = Pt(font_size.pt * 1.35)
        for run in p.runs:
            set_ea_font(run, ea_font)
    return txBox

def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    _clean_shape(shape)
    return shape

def add_hline(slide, x, y, length, color=BLACK, thickness=Pt(0.5)):
    h = max(int(thickness), Emu(6350))
    return add_rect(slide, x, y, length, h, color)

def add_oval(slide, x, y, letter, size=Inches(0.45),
             bg=NAVY, fg=WHITE):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    c.fill.solid()
    c.fill.fore_color.rgb = bg
    c.line.fill.background()
    tf = c.text_frame
    tf.paragraphs[0].text = letter
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.name = 'Arial'
    tf.paragraphs[0].font.color.rgb = fg
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    for run in tf.paragraphs[0].runs:
        set_ea_font(run, 'KaiTi')
    bodyPr = tf._txBody.find(qn('a:bodyPr'))
    bodyPr.set('anchor', 'ctr')
    _clean_shape(c)
    return c

def add_action_title(slide, text, title_size=TITLE_SIZE):
    add_text(slide, Inches(0.8), Inches(0.15), Inches(11.7), Inches(0.9),
             text, font_size=title_size, font_color=BLACK, bold=True,
             font_name='Georgia', ea_font='KaiTi', anchor=MSO_ANCHOR.MIDDLE)
    add_hline(slide, Inches(0.8), Inches(1.05), Inches(11.7),
             color=BLACK, thickness=Pt(0.5))

def add_source(slide, text, y=Inches(7.05)):
    add_text(slide, Inches(0.8), y, Inches(11), Inches(0.3),
             text, font_size=Pt(9), font_color=MED_GRAY)

def add_page_number(slide, num, total):
    add_text(slide, Inches(12.2), Inches(7.1), Inches(1), Inches(0.3),
             f"{num}/{total}", font_size=Pt(9), font_color=MED_GRAY,
             alignment=PP_ALIGN.RIGHT)

def full_cleanup(outpath):
    tmppath = outpath + '.tmp'
    with zipfile.ZipFile(outpath, 'r') as zin:
        with zipfile.ZipFile(tmppath, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename.endswith('.xml'):
                    root = etree.fromstring(data)
                    ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
                    ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
                    for style in root.findall(f'.//{{{ns_p}}}style'):
                        style.getparent().remove(style)
                    if 'theme' in item.filename.lower():
                        for tag in ['outerShdw', 'innerShdw', 'scene3d', 'sp3d']:
                            for el in root.findall(f'.//{{{ns_a}}}{tag}'):
                                el.getparent().remove(el)
                    data = etree.tostring(root, xml_declaration=True,
                                          encoding='UTF-8', standalone=True)
                zout.writestr(item, data)
    os.replace(tmppath, outpath)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BL = prs.slide_layouts[6]

total_slides = 10

# Slide 1: Cover
s1 = prs.slides.add_slide(BL)
add_rect(s1, 0, 0, SW, Inches(0.05), NAVY)
add_text(s1, Inches(1), Inches(2.2), Inches(11), Inches(1.0),
         '用AI赋能软件开发', font_size=Pt(44), font_name='Georgia',
         font_color=NAVY, bold=True, ea_font='KaiTi')
add_text(s1, Inches(1), Inches(3.5), Inches(11), Inches(0.6),
         '获取管理层支持的战略方案', font_size=Pt(24),
         font_color=DARK_GRAY, ea_font='KaiTi')
add_text(s1, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
         '演示文稿  |  2026年3月', font_size=BODY_SIZE,
         font_color=MED_GRAY, ea_font='KaiTi')
add_hline(s1, Inches(1), Inches(6.8), Inches(4), color=NAVY, thickness=Pt(2))
add_page_number(s1, 1, total_slides)

# Slide 2: Table of Contents
s2 = prs.slides.add_slide(BL)
add_action_title(s2, '目录')
items = [
    ('1', '背景与挑战', '当前软件开发面临的困境'),
    ('2', 'AI赋能开发', 'AI如何改变软件开发流程'),
    ('3', '核心价值', 'AI开发带来的业务价值'),
    ('4', '实施路径', '落地计划与时间表'),
    ('5', '投资回报', '成本收益分析'),
    ('6', '下一步', '需要决策的事项')
]
iy = Inches(1.6)
for num, title, desc in items:
    add_oval(s2, LM, iy, num, size=Inches(0.5))
    add_text(s2, LM + Inches(0.7), iy, Inches(4.0), Inches(0.4),
             title, font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True)
    add_text(s2, Inches(5.5), iy + Inches(0.05), Inches(6.5), Inches(0.4),
             desc, font_size=BODY_SIZE, font_color=MED_GRAY)
    iy += Inches(0.7)
    add_hline(s2, LM, iy, CONTENT_W, LINE_GRAY)
    iy += Inches(0.3)
add_page_number(s2, 2, total_slides)

# Slide 3: Executive Summary
s3 = prs.slides.add_slide(BL)
add_action_title(s3, '执行摘要')
add_rect(s3, LM, Inches(1.4), CONTENT_W, Inches(1.0), NAVY)
add_text(s3, LM + Inches(0.3), Inches(1.4), CONTENT_W - Inches(0.6), Inches(1.0),
         'AI是软件工程未来的必然趋势，早期布局将获得显著竞争优势',
         font_size=SUB_HEADER_SIZE, font_color=WHITE, bold=True,
         anchor=MSO_ANCHOR.MIDDLE)
points = [
    ('1', '效率提升', '开发效率预计提升40-60%'),
    ('2', '成本优化', '年度研发成本降低30%'),
    ('3', '质量保障', '代码缺陷率下降50%以上')
]
iy = Inches(2.8)
for num, title, desc in points:
    add_oval(s3, LM, iy, num)
    add_text(s3, LM + Inches(0.6), iy, Inches(3.5), Inches(0.4),
             title, font_size=BODY_SIZE, font_color=NAVY, bold=True)
    add_text(s3, Inches(5.0), iy, Inches(7.5), Inches(0.4),
             desc, font_size=BODY_SIZE)
    iy += Inches(0.6)
    add_hline(s3, LM, iy, CONTENT_W, LINE_GRAY)
    iy += Inches(0.3)
add_source(s3, '基于行业基准分析与内部评估')
add_page_number(s3, 3, total_slides)

# Slide 4: Background & Challenges
s4 = prs.slides.add_slide(BL)
add_action_title(s4, '背景与挑战：软件开发困境')
# Left column
add_text(s4, LM, Inches(1.5), Inches(5.5), Inches(0.4),
         '当前痛点', font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True)
add_hline(s4, LM, Inches(2.0), Inches(5.5), BLACK, Pt(0.5))
add_text(s4, LM, Inches(2.2), Inches(5.5), Inches(3.5),
         ['• 人力成本持续上涨', '', '• 需求变更频繁', '', '• 技术债务累积', '', '• 人才招聘困难', '', '• 测试周期冗长'],
         line_spacing=Pt(8))
# Right column
add_text(s4, Inches(6.8), Inches(1.5), Inches(5.5), Inches(0.4),
         '行业趋势', font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True)
add_hline(s4, Inches(6.8), Inches(2.0), Inches(5.5), BLACK, Pt(0.5))
add_text(s4, Inches(6.8), Inches(2.2), Inches(5.5), Inches(3.5),
         ['• 全球AI开发工具市场增长迅速', '', '• 头部科技企业已全面部署', '', '• 开源生态日益成熟', '', '• 客户对交付速度要求提高'],
         line_spacing=Pt(8))
add_source(s4, '来源：2026年软件开发行业报告')
add_page_number(s4, 4, total_slides)

# Slide 5: Big Number
s5 = prs.slides.add_slide(BL)
add_action_title(s5, 'AI赋能开发：效率飞跃')
add_rect(s5, LM, Inches(1.4), Inches(3.5), Inches(1.8), NAVY)
add_text(s5, LM + Inches(0.2), Inches(1.5), Inches(3.1), Inches(0.8),
         '60%', font_size=Pt(44), font_color=WHITE, bold=True,
         font_name='Georgia', alignment=PP_ALIGN.CENTER)
add_text(s5, LM + Inches(0.2), Inches(2.3), Inches(3.1), Inches(0.7),
         '开发效率提升', font_size=Pt(12), font_color=WHITE, alignment=PP_ALIGN.CENTER)
add_text(s5, Inches(5.0), Inches(1.5), Inches(7.5), Inches(2.0),
         '通过AI辅助编程、自动测试生成、\n智能代码审查等技术手段，\n显著缩短交付周期',
         font_size=BODY_SIZE, line_spacing=Pt(8))
# Three stats
stats = [('40%', '代码编写时间减少'), ('50%', '缺陷发现提前'), ('70%', '文档自动生成')]
sw = Inches(3.5)
sg = (CONTENT_W - sw * 3) / 2
for i, (big, small) in enumerate(stats):
    sx = LM + (sw + sg) * i
    add_rect(s5, sx, Inches(4.0), sw, Inches(1.8), BG_GRAY)
    add_text(s5, sx + Inches(0.2), Inches(4.1), sw - Inches(0.4), Inches(0.7),
             big, font_size=Pt(28), font_color=NAVY, bold=True,
             font_name='Georgia', alignment=PP_ALIGN.CENTER)
    add_text(s5, sx + Inches(0.2), Inches(4.85), sw - Inches(0.4), Inches(0.6),
             small, font_size=BODY_SIZE, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
add_source(s5, '来源：GitHub Copilot企业版效能报告')
add_page_number(s5, 5, total_slides)

# Slide 6: Three Pillars
s6 = prs.slides.add_slide(BL)
add_action_title(s6, 'AI开发的核心价值')
pillars = [
    ('效率提升', ['智能代码补全', '自动化测试生成', '快速原型构建']),
    ('质量保障', ['代码静态分析', '漏洞自动检测', '性能优化建议']),
    ('知识传承', ['文档自动生成', '代码注释智能化', '经验沉淀复用'])
]
pw = Inches(3.5)
pg = (CONTENT_W - pw * 3) / 2
for i, (title, points) in enumerate(pillars):
    px = LM + (pw + pg) * i
    add_rect(s6, px, Inches(1.5), pw, Inches(0.6), NAVY)
    add_text(s6, px + Inches(0.15), Inches(1.5), pw - Inches(0.3), Inches(0.6),
             title, font_size=SUB_HEADER_SIZE, font_color=WHITE, bold=True,
             anchor=MSO_ANCHOR.MIDDLE, alignment=PP_ALIGN.CENTER)
    add_rect(s6, px, Inches(2.1), pw, Inches(4.0), BG_GRAY)
    add_text(s6, px + Inches(0.2), Inches(2.3), pw - Inches(0.4), Inches(3.5),
             [f'• {p}' for p in points], line_spacing=Pt(10))
add_source(s6, '来源：内部技术评估')
add_page_number(s6, 6, total_slides)

# Slide 7: Timeline
s7 = prs.slides.add_slide(BL)
add_action_title(s7, '实施路径：分阶段推进')
# Timeline bar
add_hline(s7, LM + Inches(0.5), Inches(3.0), Inches(10.7), LINE_GRAY, Pt(2))
milestones = [
    ('Q2 2026', '试点启动', '选择1-2个核心项目\n引入AI开发工具'),
    ('Q3 2026', '扩大应用', '推广至50%研发团队\n建立最佳实践'),
    ('Q4 2026', '全面落地', '全员覆盖\n形成标准流程'),
    ('2027', '持续优化', '效果评估\n迭代升级')
]
spacing = Inches(10.7) / (len(milestones) - 1)
for i, (label, title, desc) in enumerate(milestones):
    mx = LM + Inches(0.5) + spacing * i
    add_oval(s7, mx - Inches(0.225), Inches(2.775), str(i + 1))
    add_text(s7, mx - Inches(1.0), Inches(2.0), Inches(2.0), Inches(0.5),
             label, font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_text(s7, mx - Inches(1.0), Inches(3.5), Inches(2.0), Inches(1.5),
             title, font_size=BODY_SIZE, font_color=NAVY, bold=True,
             alignment=PP_ALIGN.CENTER)
    add_text(s7, mx - Inches(1.0), Inches(4.7), Inches(2.0), Inches(1.5),
             desc, font_size=Pt(10), alignment=PP_ALIGN.CENTER)
add_source(s7, '来源：技术规划路线图')
add_page_number(s7, 7, total_slides)

# Slide 8: ROI
s8 = prs.slides.add_slide(BL)
add_action_title(s8, '投资回报分析')
# Big numbers
add_rect(s8, LM, Inches(1.4), Inches(4.0), Inches(2.0), NAVY)
add_text(s8, LM + Inches(0.2), Inches(1.5), Inches(3.6), Inches(0.8),
         '300万', font_size=Pt(40), font_color=WHITE, bold=True,
         font_name='Georgia', alignment=PP_ALIGN.CENTER)
add_text(s8, LM + Inches(0.2), Inches(2.4), Inches(3.6), Inches(0.7),
         '年度研发成本节省', font_size=Pt(12), font_color=WHITE, alignment=PP_ALIGN.CENTER)

add_rect(s8, Inches(5.0), Inches(1.4), Inches(4.0), Inches(2.0), BG_GRAY)
add_text(s8, Inches(5.2), Inches(1.5), Inches(3.6), Inches(0.8),
         '1.5年', font_size=Pt(40), font_color=NAVY, bold=True,
         font_name='Georgia', alignment=PP_ALIGN.CENTER)
add_text(s8, Inches(5.2), Inches(2.4), Inches(3.6), Inches(0.7),
         '投资回收期', font_size=Pt(12), font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

add_rect(s8, Inches(9.5), Inches(1.4), Inches(3.5), Inches(2.0), BG_GRAY)
add_text(s8, Inches(9.7), Inches(1.5), Inches(3.1), Inches(0.8),
         '5x', font_size=Pt(40), font_color=NAVY, bold=True,
         font_name='Georgia', alignment=PP_ALIGN.CENTER)
add_text(s8, Inches(9.7), Inches(2.4), Inches(3.1), Inches(0.7),
         '三年ROI', font_size=Pt(12), font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# Cost breakdown
add_text(s8, LM, Inches(4.0), Inches(11.7), Inches(0.4),
         '主要投资项', font_size=SUB_HEADER_SIZE, font_color=NAVY, bold=True)
add_hline(s8, LM, Inches(4.5), Inches(11.7), BLACK, Pt(0.5))
items = [
    ('AI开发工具许可', '80万/年'),
    ('培训与变革管理', '30万'),
    ('基础设施升级', '50万')
]
iy = Inches(4.7)
for item, cost in items:
    add_text(s8, LM, iy, Inches(8), Inches(0.4), item, font_size=BODY_SIZE)
    add_text(s8, Inches(9), iy, Inches(3), Inches(0.4), cost, font_size=BODY_SIZE, 
             font_color=NAVY, bold=True, alignment=PP_ALIGN.RIGHT)
    iy += Inches(0.5)
add_source(s8, '来源：财务部门估算')
add_page_number(s8, 8, total_slides)

# Slide 9: Next Steps
s9 = prs.slides.add_slide(BL)
add_action_title(s9, '下一步：需要决策的事项')
actions = [
    ('批准试点项目', 'Q2 2026启动', '选择2个核心项目进行AI开发试点', 'CTO'),
    ('预算审批', '2026年3月', '批准首期160万投资预算', 'CFO'),
    ('团队组建', '2026年4月', '成立AI开发变革小组', 'HR')
]
cw = Inches(3.5)
cg = (CONTENT_W - cw * 3) / 2
for i, (title, timeline, desc, owner) in enumerate(actions):
    cx = LM + (cw + cg) * i
    add_rect(s9, cx, Inches(1.5), cw, Inches(0.6), NAVY)
    add_text(s9, cx + Inches(0.15), Inches(1.5), cw - Inches(0.3), Inches(0.6),
             title, font_size=BODY_SIZE, font_color=WHITE, bold=True,
             anchor=MSO_ANCHOR.MIDDLE, alignment=PP_ALIGN.CENTER)
    add_rect(s9, cx, Inches(2.1), cw, Inches(0.4), BG_GRAY)
    add_text(s9, cx + Inches(0.15), Inches(2.1), cw - Inches(0.3), Inches(0.4),
             timeline, font_size=BODY_SIZE, font_color=NAVY, bold=True,
             anchor=MSO_ANCHOR.MIDDLE, alignment=PP_ALIGN.CENTER)
    add_text(s9, cx + Inches(0.15), Inches(2.7), cw - Inches(0.3), Inches(2.0),
             desc.split('\n'), line_spacing=Pt(8), alignment=PP_ALIGN.CENTER)
    add_hline(s9, cx + Inches(0.3), Inches(4.9), cw - Inches(0.6), LINE_GRAY)
    add_text(s9, cx + Inches(0.15), Inches(5.1), cw - Inches(0.3), Inches(0.4),
             f'负责人：{owner}', font_size=BODY_SIZE, font_color=MED_GRAY,
             alignment=PP_ALIGN.CENTER)
add_source(s9, '')
add_page_number(s9, 9, total_slides)

# Slide 10: Closing
s10 = prs.slides.add_slide(BL)
add_rect(s10, 0, 0, SW, Inches(0.05), NAVY)
add_text(s10, Inches(1.5), Inches(2.0), Inches(10.3), Inches(1.0),
         'AI赋能开发，赢得未来竞争', font_size=Pt(28), font_color=NAVY, bold=True,
         font_name='Georgia', alignment=PP_ALIGN.CENTER)
add_hline(s10, Inches(5.5), Inches(3.3), Inches(2.3), NAVY, Pt(1.5))
add_text(s10, Inches(1.5), Inches(3.8), Inches(10.3), Inches(2.0),
         '立即行动，把握AI带来的竞争优势',
         font_size=SUB_HEADER_SIZE, font_color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
add_hline(s10, Inches(1), Inches(6.8), Inches(3), NAVY, Pt(2))
add_text(s10, Inches(1.5), Inches(7.0), Inches(10.3), Inches(0.4),
         '感谢聆听  |  欢迎提问', font_size=BODY_SIZE, font_color=MED_GRAY, 
         alignment=PP_ALIGN.CENTER)
add_page_number(s10, 10, total_slides)

# Save with cleanup
outpath = 'AI-Software-Development-Proposal.pptx'
prs.save(outpath)
full_cleanup(outpath)
print(f"Presentation saved: {outpath}")
